"""PowerPoint (.pptx) document parser using python-pptx.

Extracts all text from slides, preserving slide structure with page separators.
"""
from pptx import Presentation


def extract_text_from_pptx(file_path: str) -> str:
    """Extract all text from a PowerPoint file.

    Traverses every slide and collects text from all shapes (text boxes,
    tables, etc.), inserting a slide-boundary separator after each slide
    to preserve structural context for downstream chunking.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        All extracted text joined with slide separators.
    """
    prs = Presentation(file_path)

    slides_text = []
    total_slides = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    runs = [r.text for r in paragraph.runs if r.text and r.text.strip()]
                    if runs:
                        slide_parts.append("".join(runs))

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_cells = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_cells.append(cell_text)
                    if row_cells:
                        slide_parts.append(" | ".join(row_cells))

        slide_content = "\n".join(slide_parts)

        sep = f"\n--- [第 {slide_idx}/{total_slides} 页幻灯片] ---\n"
        slides_text.append(slide_content + sep)

    return "".join(slides_text)
